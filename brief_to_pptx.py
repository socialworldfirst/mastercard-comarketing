#!/usr/bin/env python3
"""Render brief.html into an editable two-slide .pptx on the WF template.

Bespoke to the brief's layout (lead + film references, three columns). Everything
is native: text frames, pictures with live hyperlinks, no screenshots.

    python3 brief_to_pptx.py            # -> ~/Downloads/mastercard-comarketing-brief.pptx
"""
import re, html, pathlib
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

HERE = pathlib.Path(__file__).parent
TEMPLATE = pathlib.Path.home() / "Documents/Claude/Assets/WF_PPT_Template/WF_PPT_Template_Intl_Apr2025.pptx"
OUT = pathlib.Path.home() / "Downloads" / "mastercard-comarketing-brief.pptx"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Poppins"
raw = (HERE / "brief.html").read_text(encoding="utf-8")


def text_of(f):
    f = re.sub(r"<[^>]+>", "", f)
    return html.unescape(re.sub(r"\s+", " ", f)).strip()


# ---------- parse the brief ------------------------------------------------
title = text_of(re.search(r'<h1 class="title"[^>]*>(.*?)</h1>', raw, re.S).group(1))
lead = [text_of(p) for p in re.findall(r'<div class="lead">(.*?)</div>', raw, re.S)[0]
        .split("</p>") if text_of(p)]
films = re.findall(r'<a class="film" href="([^"]+)"[^>]*>\s*<div class="shot">'
                   r'<img src="([^"]+)"[^>]*></div>\s*</a>\s*<span>(.*?)</span>', raw, re.S)
if not films:  # caption inside the anchor in the current markup
    films = [(h, i, text_of(c)) for h, i, c in re.findall(
        r'<a class="film" href="([^"]+)"[^>]*>.*?<img src="([^"]+)".*?<span>(.*?)</span>', raw, re.S)]
p2_title = text_of(re.search(r'<h2 class="title">(.*?)</h2>', raw, re.S).group(1))
cols = []
for c in re.findall(r'<div class="col">(.*?)</div>', raw, re.S):
    head = text_of(re.search(r"<h3[^>]*>(.*?)</h3>", c, re.S).group(1))
    items = [text_of(x) for x in re.findall(r"<li>(.*?)</li>", c, re.S)]
    cols.append((head, items))
more = re.search(r'<p class="more">(.*?)</p>', raw, re.S)
more_txt = text_of(more.group(1)) if more else ""
more_url = (re.search(r'href="([^"]+)"', more.group(1)).group(1) if more else "")

# ---------- build ----------------------------------------------------------
prs = Presentation(str(TEMPLATE))
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    prs.part.drop_rel(sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
    sldIdLst.remove(sldId)
layouts = {l.name: l for m in prs.slide_masters for l in m.slide_layouts}


def add(layout):
    slide = prs.slides.add_slide(layouts[layout])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, lines, size, bold=False, space=6, bullet=False, line=1.4, alpha=None,
          caps=False, url=None):
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        p.line_spacing = line
        r = p.add_run()
        r.text = ("•  " if bullet else "") + (txt.upper() if caps else txt)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = WHITE
        if url:
            r.hyperlink.address = url
        if alpha is not None:
            a = OxmlElement("a:alpha"); a.set("val", str(alpha))
            r.font.color._xFill.srgbClr.append(a)


def rule(slide, x, y, w):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(0.75))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    a = OxmlElement("a:alpha"); a.set("val", "30000")
    sh.fill.fore_color._xFill.srgbClr.append(a)
    sh.line.fill.background(); sh.shadow.inherit = False


# --- slide 1 ---------------------------------------------------------------
s1 = add("Cover Slide (Dark)")
write(box(s1, 37.7, 96, 600, 100), [title], 40, bold=True, line=1.1)
write(box(s1, 37.7, 216, 560, 220), lead, 17, space=12, line=1.5)
write(box(s1, 670.3, 120, 252, 16), ["References"], 10, bold=True, caps=True, alpha=70000)
fy = 146.0
for href, img, cap in films[:2]:
    src = HERE / img
    png = src.with_suffix(".png")
    if not png.exists():
        Image.open(src).convert("RGB").save(png)
    pic = s1.shapes.add_picture(str(png), Pt(670.3), Pt(fy), Pt(252), Pt(141.75))
    pic.click_action.hyperlink.address = href
    write(box(s1, 670.3, fy + 147, 252, 18), [cap], 9.5, alpha=85000, line=1.3)
    fy += 141.75 + 34

# --- slide 2 ---------------------------------------------------------------
s2 = add("Text 1-Column (Light)")   # the pink content layout
write(box(s2, 37.7, 27.6, 768, 40), [p2_title], 32, bold=True, line=1.1)
colw, gap = 270.9, 36.0
for i, (head, items) in enumerate(cols):
    x = 37.7 + i * (colw + gap)
    write(box(s2, x, 110, colw, 22), [head], 15, bold=True, line=1.2)
    rule(s2, x, 138, colw)
    write(box(s2, x, 150, colw, 300), items, 12.5, space=9, bullet=True, line=1.45)
if more_txt:
    # plain text, no hyperlink run: PowerPoint forces theme blue on linked runs
    tf = box(s2, 37.7, 478, 700, 20)
    write(tf, [more_txt], 10.5, alpha=85000)

prs.save(str(OUT))
print(f"2 slides -> {OUT}")
