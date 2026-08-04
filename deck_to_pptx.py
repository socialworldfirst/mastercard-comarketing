#!/usr/bin/env python3
"""Render a deck's index.src.html into an editable .pptx on the WF template.

Everything is native PowerPoint: text frames, tables and shapes, never a picture
of a slide. Backgrounds come from the real template layouts, so the output opens
on-brand and every word stays editable.

    python3 deck_to_pptx.py                 # -> ~/Downloads/<slug>.pptx
    python3 deck_to_pptx.py --out path.pptx

Run from a deck folder (needs index.src.html alongside).
"""
import re, sys, html, pathlib, copy
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement

HERE = pathlib.Path.cwd()
SRC = HERE / "index.src.html"
TEMPLATE = pathlib.Path.home() / "Documents/Claude/Assets/WF_PPT_Template/WF_PPT_Template_Intl_Apr2025.pptx"
if not SRC.exists():
    raise SystemExit("no index.src.html here — run this from a deck folder")
if not TEMPLATE.exists():
    raise SystemExit(f"template not found: {TEMPLATE}")

raw = SRC.read_text(encoding="utf-8")
m = re.search(r"<title>(.*?)</title>", raw, re.S)
DECK_TITLE = html.unescape(m.group(1).strip()) if m else HERE.name

# ---------- design tokens (theme "WF 17") ---------------------------------
PINK  = RGBColor(0xFF, 0x00, 0x51)
INK   = RGBColor(0x1F, 0x23, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xEC, 0xED, 0xEE)
PURPLE     = RGBColor(0x32, 0x00, 0x6D)
PURPLE_MID = RGBColor(0x8F, 0x66, 0xC7)
LILAC      = RGBColor(0xD8, 0xC9, 0xE8)
WARMGREY   = RGBColor(0x90, 0x96, 0x9A)
GREY       = RGBColor(0xD5, 0xD5, 0xD5)
BAR_FILLS = [PURPLE, PURPLE_MID, LILAC, WARMGREY, GREY]
BAR_INK   = [WHITE, WHITE, INK, WHITE, INK]

FONT = "Poppins"
MARGIN_L, TITLE_T, TITLE_W = 37.7, 27.6, 828.0

# deck background class -> template layout name
LAYOUT_FOR = {
    "bg-cover":           "Cover Slide (Dark)",
    "bg-cover-ship":      "1_Cover Slide (Dark)",
    "bg-cover-sme":       "2_Cover Slide (Dark)",
    "bg-pink":            "Text 1-Column (Light)",
    "bg-light":           "1_Text 1-Column (Light)",
    "bg-purple":          "2_Text 1-Column (Purple)",
    "bg-section":         "1_Section Divider (Dark)",
    "bg-section-light":   "2_Section Divider (Light)",
    "bg-section-purple":  "2_Section Divider (Purple)",
    "bg-quote":           "Quotes (Light)",
    "bg-thanks":          "Thank you",
}
LIGHT_BG = {"bg-light", "bg-section-light"}


def text_of(frag):
    frag = re.sub(r"<br\s*/?>", " ", frag)
    frag = re.sub(r"<[^>]+>", "", frag)
    return html.unescape(re.sub(r"\s+", " ", frag)).strip()


def divs(frag, cls_re):
    out = []
    for m in re.finditer(r'<div class="([^"]*)"[^>]*>', frag):
        if not re.search(cls_re, m.group(1)):
            continue
        depth, i = 1, m.end()
        while depth and i < len(frag):
            nxt = re.search(r"</?div\b", frag[i:])
            if not nxt:
                break
            j = i + nxt.start()
            depth += -1 if frag[j:j + 2] == "</" else 1
            i = j + (6 if frag[j:j + 2] == "</" else 4)
        out.append(frag[m.end():i - 6])
    return out


def parse():
    slides = []
    for m in re.finditer(r'<section class="slide([^"]*)">(.*?)</section>', raw, re.S):
        cls, inner = m.group(1).strip(), m.group(2)
        s = {"cls": cls.split(), "title": None, "sub": None, "body": [], "bullets": [],
             "tables": [], "agenda": [], "refs": [], "timeline": None, "note": None,
             "quote": None, "quotebox": None}
        t = re.search(r'<h[12] class="title">(.*?)</h[12]>', inner, re.S)
        if t: s["title"] = text_of(t.group(1))
        t = re.search(r'<p class="sub">(.*?)</p>', inner, re.S)
        if t: s["sub"] = text_of(t.group(1))
        for b in divs(inner, r"\bbody\b"):
            li = re.findall(r"<li>(.*?)</li>", b, re.S)
            if li: s["bullets"] += [text_of(x) for x in li]
            else:  s["body"] += [text_of(p) for p in re.findall(r"<p>(.*?)</p>", b, re.S)]
        for a in divs(inner, r"\bagenda\b"):
            s["agenda"] = [(text_of(x), text_of(y))
                           for x, y in re.findall(r"<b>(.*?)</b>\s*<span>(.*?)</span>", a, re.S)]
        for tb in divs(inner, r"\btbl\b"):
            rows = []
            for tr in re.finditer(r'<div class="tr([^"]*)">(.*?)</div>', tb, re.S):
                cells = [text_of(c) for c in re.findall(r"<span[^>]*>(.*?)</span>", tr.group(2), re.S)]
                if cells: rows.append((cells, "th" in tr.group(1)))
            if rows: s["tables"].append(rows)
        for rf in divs(inner, r"\brefs\b"):
            imgs = re.findall(r'<img src="([^"]+)"', rf)
            caps = re.findall(r"<h3>(.*?)</h3>\s*<p>(.*?)</p>", rf, re.S)
            s["refs"] = [(imgs[i] if i < len(imgs) else None, text_of(a), text_of(b))
                         for i, (a, b) in enumerate(caps)]
        for tl in divs(inner, r"\btimeline\b"):
            heads = [text_of(x) for x in re.findall(r"<div>(.*?)</div>", tl, re.S)][:5]
            bars = []
            for bm in re.finditer(r'<div class="bar (p\d)"[^>]*style="([^"]*)"[^>]*>(.*?)</div>', tl, re.S):
                st = bm.group(2)
                w = float(re.search(r"width:(\d+(?:\.\d+)?)px", st).group(1))
                ml = re.search(r"margin-left:(\d+(?:\.\d+)?)px", st)
                bars.append((int(bm.group(1)[1]) - 1, float(ml.group(1)) if ml else 0.0,
                             w, text_of(bm.group(3))))
            s["timeline"] = (heads, bars)
        t = re.search(r'<p class="note">(.*?)</p>', inner, re.S)
        if t: s["note"] = text_of(t.group(1))
        t = re.search(r'<p class="quote[^"]*">(.*?)</p>', inner, re.S)
        if t: s["quote"] = text_of(t.group(1))
        t = re.search(r'<div class="quotebox"[^>]*>(.*?)</div>', inner, re.S)
        if t: s["quotebox"] = text_of(t.group(1))
        slides.append(s)
    return slides


# ---------- presentation --------------------------------------------------
prs = Presentation(str(TEMPLATE))

# strip the template's 32 sample slides, keep masters and layouts
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

layouts = {l.name: l for m in prs.slide_masters for l in m.slide_layouts}


def add(cls):
    bg = next((c for c in cls if c in LAYOUT_FOR), "bg-pink")
    slide = prs.slides.add_slide(layouts[LAYOUT_FOR[bg]])
    # remove the layout's prompt placeholders; we position our own boxes
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide, (INK if bg in LIGHT_BG else WHITE)


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, lines, size, colour, bold=False, space=6, bullet=False, line=1.35):
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        p.line_spacing = line
        r = p.add_run()
        r.text = ("•  " if bullet else "") + txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = colour


def _no_borders(cell, rule_colour=None):
    """PowerPoint's default table style paints borders and banding. Strip them,
    then optionally draw a single hairline under the row to match the HTML."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = OxmlElement(tag)
        ln.set("w", "9525")
        if tag == "a:lnB" and rule_colour is not None:
            fill = OxmlElement("a:solidFill")
            clr = OxmlElement("a:srgbClr")
            clr.set("val", str(rule_colour))
            alpha = OxmlElement("a:alpha"); alpha.set("val", "26000")
            clr.append(alpha); fill.append(clr); ln.append(fill)
        else:
            ln.append(OxmlElement("a:noFill"))
        tcPr.append(ln)


def table(slide, rows, x, y, w, ink):
    ncols = max(len(c) for c, _ in rows)
    nrows = len(rows)
    gt = slide.shapes.add_table(nrows, ncols, Pt(x), Pt(y), Pt(w), Pt(23 * nrows)).table
    gt.first_row = False          # no coloured header band
    gt.horz_banding = False       # no zebra rows
    first = 176 if w > 700 else 132
    gt.columns[0].width = Pt(first)
    rest = (w - first) / max(1, ncols - 1)
    for c in range(1, ncols):
        gt.columns[c].width = Pt(rest)
    for ri, (cells, is_head) in enumerate(rows):
        gt.rows[ri].height = Pt(23)
        for ci in range(ncols):
            cell = gt.cell(ri, ci)
            cell.fill.background()
            _no_borders(cell, ink)
            cell.margin_left = cell.margin_right = Pt(0)
            cell.margin_top = cell.margin_bottom = Pt(4)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = cells[ci] if ci < len(cells) else ""
            r.font.size = Pt(10 if is_head else 11)
            r.font.name = FONT
            r.font.bold = (ci == 0 and not is_head)
            r.font.color.rgb = ink        # follows the slide, not always dark
            if is_head:
                r.font.color.rgb = ink
                solidFill = r.font.color._xFill
                alpha = OxmlElement("a:alpha"); alpha.set("val", "62000")
                solidFill.srgbClr.append(alpha)
    return gt


def build():
    slides = parse()
    for n, s in enumerate(slides, 1):
        cls = s["cls"]
        slide, ink = add(cls)
        is_cover = "is-cover" in cls
        is_section = "is-section" in cls
        is_thanks = "is-thanks" in cls

        if s["title"]:
            if is_thanks:
                tf = box(slide, 0, 246, 960, 60)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                write(tf, [s["title"]], 40, ink, bold=True, line=1.0)
            elif is_cover:
                write(box(slide, MARGIN_L, 177, 546, 60), [s["title"]], 40, ink, bold=True, line=1.1)
            elif is_section:
                write(box(slide, MARGIN_L, 181, 498, 54), [s["title"]], 32, ink, bold=True, line=1.1)
            else:
                write(box(slide, MARGIN_L, TITLE_T, TITLE_W - 60, 40), [s["title"]], 32, ink, bold=True, line=1.1)

        if s["sub"]:
            y, sz = (233, 24) if is_cover else (227, 18)
            write(box(slide, MARGIN_L, y, 500, 40), [s["sub"]], sz, ink, line=1.25)

        y = 104.5
        if s["body"]:
            wide = 786 if ("wide" in raw or True) else 551.6
            h = 20 + 24 * sum(max(1, len(t) // 95 + 1) for t in s["body"])
            write(box(slide, MARGIN_L, y, wide, h), s["body"], 16, ink, space=10, line=1.4)
            y += h + 10

        if s["bullets"]:
            write(box(slide, MARGIN_L, y, 786, 30 * len(s["bullets"])),
                  s["bullets"], 16, ink, space=9, bullet=True, line=1.35)

        if s["agenda"]:
            tf = box(slide, 471.5, 185, 300, 260)
            for i, (a, b) in enumerate(s["agenda"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(4)
                r = p.add_run(); r.text = a
                r.font.size = Pt(28); r.font.name = FONT; r.font.color.rgb = ink
                q = tf.add_paragraph(); q.space_after = Pt(20)
                r = q.add_run(); r.text = b
                r.font.size = Pt(12); r.font.name = FONT; r.font.color.rgb = ink

        for ti, rows in enumerate(s["tables"]):
            table(slide, rows, MARGIN_L, y if ti == 0 else y + 190, 884.6, ink)

        if s["refs"]:
            w, gap = 430, 24
            for i, (img, head, cap) in enumerate(s["refs"][:2]):
                x = MARGIN_L + i * (w + gap)
                if img and (HERE / img).exists():
                    png = (HERE / img).with_suffix(".png")
                    if not png.exists():
                        from PIL import Image
                        Image.open(HERE / img).convert("RGB").save(png)
                    slide.shapes.add_picture(str(png), Pt(x), Pt(150), Pt(w), Pt(w * 9 / 16))
                write(box(slide, x, 150 + w * 9 / 16 + 10, w, 24), [head], 16, ink, bold=True, line=1.2)
                write(box(slide, x, 150 + w * 9 / 16 + 34, w, 60), [cap], 11, ink, line=1.45)

        if s["timeline"]:
            heads, bars = s["timeline"]
            tlx, tly, tlw, tlh = 48.4, 176.4, 800.0, 271.7
            colw = tlw / max(1, len(heads))
            for i, hd in enumerate(heads):
                ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(tlx + i * colw), Pt(tly), Pt(0.75), Pt(tlh))
                ln.fill.solid(); ln.fill.fore_color.rgb = ink; ln.line.fill.background()
                ln.shadow.inherit = False
                write(box(slide, tlx + i * colw + 10, tly - 2, colw - 12, 20), [hd], 14, ink, line=1.1)
            for bi, (pi, ml, bw, label) in enumerate(bars):
                sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Pt(tlx + 2 + ml), Pt(tly + 40 + bi * 41), Pt(bw), Pt(36.4))
                sh.adjustments[0] = 0.13
                sh.fill.solid(); sh.fill.fore_color.rgb = BAR_FILLS[pi]
                sh.line.color.rgb = WHITE; sh.line.width = Pt(2.25)
                sh.shadow.inherit = False
                p = sh.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                sh.text_frame.margin_left = Pt(14)
                r = p.add_run(); r.text = label
                r.font.size = Pt(14); r.font.name = FONT; r.font.color.rgb = BAR_INK[pi]

        if s["quote"]:
            write(box(slide, MARGIN_L, 95.7, 706, 324), [s["quote"]], 32, ink, line=1.18)

        if s["quotebox"]:
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(MARGIN_L), Pt(210), Pt(884.6), Pt(100))
            panel.adjustments[0] = 0.03
            panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
            panel.line.fill.background(); panel.shadow.inherit = False
            tf = panel.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Pt(20); tf.margin_top = Pt(14)
            write(tf, [s["quotebox"]], 12, INK, line=1.5)

        if s["note"]:
            write(box(slide, MARGIN_L, 470, 700, 30), [s["note"]], 10, ink, line=1.4)

        # No page number here: the template master already draws an automatic one
        # at 856.0, 491.1. Adding a second box put two numbers on top of each other.

    return len(slides)


count = build()
out = pathlib.Path(sys.argv[sys.argv.index("--out") + 1]) if "--out" in sys.argv \
      else pathlib.Path.home() / "Downloads" / f"{HERE.name}.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"{count} slides -> {out}")
